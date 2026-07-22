#!/usr/bin/env python3
"""WS3: extract an existing .pptx into structured deck content for restyling.

We do NOT reproduce the original PPT's styling (no fidelity mode — our value is the
professional RESTYLE via an html-deck design system). We extract the CONTENT:
per-slide {title, bullets (with level), images KEPT, speaker notes}, and emit a
content JSON that track-e2 reflows into the chosen design system.

Images are PRESERVED (they're professionally made — discarding them changes the
meaning). They are exported to an assets dir next to the JSON, referenced by a
workspace-relative path. Decorative/chrome text (page numbers, footers, empty
placeholders) is dropped by heuristic; meaningful text is kept.

Usage: pptx_to_deck.py <input.pptx> [--out <dir>] [--json]
Output: <out>/deck_content.json  +  <out>/images/slideNN_img.<ext>
"""
import sys, os, json, argparse, re
from pathlib import Path

# heuristics for dropping decorative/chrome text (not meaningful content)
_CHROME_RE = re.compile(r'^\s*(\d{1,3}|page \d+|\d+\s*/\s*\d+|©.*|confidential.*|draft.*)\s*$', re.I)

def _is_chrome(text, shape, slide_h_emu):
    """Drop page numbers, footers, copyright, and tiny bottom-strip text."""
    t = text.strip()
    if not t:
        return True
    if _CHROME_RE.match(t):
        return True
    # very short all-caps single token near the bottom → likely a footer tag
    if len(t) <= 3 and shape is not None and shape.top is not None:
        if slide_h_emu and shape.top > slide_h_emu * 0.9:
            return True
    return False

def extract(pptx_path, out_dir):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    img_dir = out_dir / "images"
    slides = []
    img_count = 0

    def collect_pictures(shapes, i):
        """Recurse into groups — pictures are often nested inside GROUP shapes."""
        nonlocal img_count
        found = []
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                found += collect_pictures(sh.shapes, i)
            elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = sh.image
                    ext = img.ext or "png"
                    img_dir.mkdir(parents=True, exist_ok=True)
                    fn = f"slide{i:02d}_img{img_count}.{ext}"
                    (img_dir / fn).write_bytes(img.blob)
                    found.append((f"images/{fn}", sh.top or 0, (sh.width or 0) * (sh.height or 0)))
                    img_count += 1
                except Exception:
                    pass
        return found

    for i, slide in enumerate(prs.slides):
        title = ""
        bullets = []          # [{"text","level"}]
        images = []           # workspace-relative paths
        notes = ""
        title_shape = slide.shapes.title if slide.shapes.title is not None else None

        # images (incl. nested in groups), KEPT
        pics = collect_pictures(slide.shapes, i)
        images = [p for p, _, _ in pics]

        # gather all meaningful text shapes (skip pictures/groups handled above)
        # each: (text, level, top, font_size_pt, is_placeholder_title)
        text_items = []
        for sh in slide.shapes:
            if sh.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP):
                continue
            if not sh.has_text_frame:
                continue
            is_title_ph = (title_shape is not None and sh == title_shape)
            top = sh.top or 0
            for para in sh.text_frame.paragraphs:
                txt = para.text.strip()
                if not txt or _is_chrome(txt, sh, slide_h):
                    continue
                # font size (pt) for title inference when no placeholder exists
                fsz = 0
                for run in para.runs:
                    if run.font.size is not None:
                        fsz = max(fsz, run.font.size.pt)
                text_items.append({
                    "text": txt, "level": int(para.level or 0),
                    "top": top, "fsz": fsz, "is_title_ph": is_title_ph,
                })

        # TITLE resolution: prefer the placeholder title; else infer — the topmost
        # largest-font text line (AI-generated decks use AUTO_SHAPE, no title ph).
        # Skip pure-emoji/symbol shapes (a big decorative 🐝 must NOT win the title).
        def _has_words(s):
            return bool(re.search(r'[0-9A-Za-z一-鿿]', s))
        title_item = None
        ph_titles = [t for t in text_items if t["is_title_ph"] and _has_words(t["text"])]
        wordy = [t for t in text_items if _has_words(t["text"])]
        if ph_titles:
            title_item = ph_titles[0]
        elif wordy:
            # infer: rank by (largest font, nearest top) among text that has words
            title_item = sorted(wordy, key=lambda t: (-t["fsz"], t["top"]))[0]
        if title_item:
            title = title_item["text"]

        for t in text_items:
            if t is title_item:
                continue
            bullets.append({"text": t["text"], "level": t["level"]})

        # speaker notes (kept — useful context for restyle)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "index": i,
            "title": title,
            "bullets": bullets,
            "images": images,
            "notes": notes,
        })

    return {
        "source": os.path.basename(pptx_path),
        "slide_count": len(slides),
        "image_count": img_count,
        "aspect": "16:9" if abs((slide_w / slide_h) - 16/9) < 0.1 else f"{slide_w}x{slide_h}",
        "slides": slides,
        "restyle_note": "Content extracted for professional restyle via an html-deck "
                        "design system. Original styling intentionally discarded; images "
                        "preserved. Feed to track-e2 with the user-selected system.",
    }

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--out", default=None)
    ap.add_argument("--json", action="store_true")
    a = ap.parse_args()

    src = Path(a.pptx).expanduser()
    if not src.is_file():
        print(json.dumps({"error": f"not found: {src}"})); sys.exit(1)
    out_dir = Path(a.out).expanduser() if a.out else src.parent / f"{src.stem}_deck_content"
    out_dir.mkdir(parents=True, exist_ok=True)

    try:
        data = extract(str(src), out_dir)
    except Exception as e:
        # a corrupt / non-pptx file (renamed .pptx) must fail with a clean JSON
        # error, not a stack trace (Gate-2 MEDIUM).
        print(json.dumps({"error": f"not a readable .pptx: {type(e).__name__}: {e}"}))
        sys.exit(1)
    (out_dir / "deck_content.json").write_text(json.dumps(data, ensure_ascii=False, indent=2))

    if a.json:
        print(json.dumps(data, ensure_ascii=False, indent=2))
    else:
        print(f"extracted {data['slide_count']} slides, {data['image_count']} images → {out_dir}")
        for s in data["slides"][:6]:
            print(f"  slide {s['index']}: title={s['title'][:40]!r} bullets={len(s['bullets'])} imgs={len(s['images'])}")

if __name__ == "__main__":
    main()
